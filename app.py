# THIS CODE AIMS AT INTEGRATING THE CORRECT IMAGE DOWNLOADING WITH THE MODIFIED PROMPT GENERATORS
import requests
import os
from flask import Flask, render_template, request, jsonify
import collections
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from bs4 import BeautifulSoup
import urllib.parse
import re

API_KEY = "KEY"
app = Flask(__name__)
slides = []
def download_images(query, num_images, i):
    query = urllib.parse.quote(query)
    url = f"https://www.google.com/search?q={query}&tbm=isch"
    response = requests.get(url)
    soup = BeautifulSoup(response.text, "html.parser")
    image_elements = soup.find_all("img")
    count = 0
    image_urls = []
    for element in image_elements:
        image_url = element["src"]
        try:
            image_response = requests.get(image_url, stream=True)
            with open(f"image_{i}.jpg", "wb") as f:
                for chunk in image_response.iter_content(1024):
                    f.write(chunk)
            image_urls.append(f"image_{i}.jpg")
            count += 1
            if count == num_images:
                break
        except requests.exceptions.RequestException as e:
            print(f"Error downloading image: {str(e)}")
    return image_urls
@app.route("/")
def index():
    title1 = request.args.get("title1")
    grade = request.args.get("grade")
    num_slides = request.args.get("num_slides")
    if title1 is None or grade is None or num_slides is None:
        return render_template("index.html")
    prompt = f"""
    Give content for {num_slides} slides on the topic '{title1}' for Grade {grade} Students 
     in following format 
'Title of Slide'
The solar system is a group of planets, moons, asteroids, comets, and other objects that revolve around a star called the Sun.
The Sun is a huge ball of hot, glowing gases, and it is at the center of our solar system.
The solar system is about 4.6 billion years old. 
    """
    print(prompt)
    response = requests.post(
        "https://api.openai.com/v1/engines/text-davinci-003/completions",
        headers={
            "Authorization": f"Bearer {API_KEY}",
            "Content-Type": "application/json",
        },
        json={
            "prompt": f"Give content for {num_slides} slides on the topic '{title1}' for Grade {grade} Students  in following format 'Title of Slide'The solar system is a group of planets, moons, asteroids, comets, and other objects that revolve around a star called the Sun.The Sun is a huge ball of hot, glowing gases, and it is at the center of our solar system.The solar system is about 4.6 billion years old.",
            "temperature": 0,
            "max_tokens": 2000,
        }
    )
    response_json = response.json()
    slides = response_json["choices"][0]["text"].split("\n")
    slide_list = slides
    # print(slide_list)
    
    slide_list_filtered = []
    for slide in slide_list:
        if slide != '':
            title = re.search(r'(Slide \d+: )?(Title: )?(.+)', slide)
            if title:
                slide_list_filtered.append(title.group(3))
            else:
                slide_list_filtered.append(slide)
    slides_data = []
    for i in range(0, len(slide_list_filtered), 2):
        title = slide_list_filtered[i]
        content = slide_list_filtered[i + 1] if i + 1 < len(slide_list_filtered) else ""
        slides_data.append((title, content))
    presentation = Presentation()
    for i, (title, content) in enumerate(slides_data):
        slide_layout = presentation.slide_layouts[1]
        slide = presentation.slides.add_slide(slide_layout)
        title_placeholder = slide.shapes.title
        title_placeholder.text = title
        content_placeholder = slide.placeholders[1]
        content_placeholder.text = content
        image_urls = download_images(title, 1, i)
        if len(image_urls) > 0:
            image_path = image_urls[0]
            image = slide.shapes.add_picture(image_path, Inches(4), Inches(1), width=Inches(4), height=Inches(4))
    presentation.save(r"C:\Users\dell\Desktop\Slide Generator\SlidesGenerator\templates\output.pptx")
    return render_template("index.html", slides=slides)
if __name__ == "__main__":
    app.run(host="192.168.56.1", port=9000, debug=True)

# THIS CODE HAS THE 'changed' PROMPT