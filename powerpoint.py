import json
import os
import requests
from flask import Flask, render_template, request, jsonify
from pptx import Presentation

API_KEY = "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"

app = Flask(__name__)
slides = []

@app.route("/")
def index():
    # Get Input from HTML File
    title = request.args.get("title")
    grade = request.args.get("grade")
    num_slides = request.args.get("num_slides")
    # Empty Check
    if title is None or grade is None or num_slides is None:
        return render_template("index.html")
    # Prompt
    prompt = f"""
    Prepare {num_slides} slides about {title} for grade {grade} students. You must be following the structure of the example given below, where each heading must contain 4 bullet points. Moreover, return response in Json format.

Example: Title: "Unbalanced trees"
1. If a BST is not balanced, its time complexity can degrade to O(n)
2. Inefficient operations: BSTs do not provide efficient support for some operations such as finding the kth smallest 
3. Not efficient on duplicates and skewness
4. Space complexity: BSTs require a lot of memory allocation
    """
    print()
    print("------------")
    print(prompt)  # Debugger Check
    print()
    print("------------")
    # ChatGPT Response
    response = requests.post(
        "https://api.openai.com/v1/engines/davinci/completions",
        headers={
            "Authorization": f"Bearer {API_KEY}",
            "Content-Type": "application/json",
        },
        data=json.dumps(
            {
                "prompt": f"Prepare {num_slides} slides about {title} for {grade} grade students.",
                "temperature": 0,
                "max_tokens": 2000,
            }
        ),
    )

    # Parse
    response_json = response.json()
    print(response_json)
    slides = response_json["choices"][0]["text"].split("\n")

    # Renders the HTML File again
    return render_template("index.html", slides=slides)


@app.route("/generate", methods=["GET"])
def generate_slides():
    title = request.args.get("title")
    num_slides = request.args.get("num_slides")

    # Create a new PowerPoint presentation
    presentation = Presentation()

    # Iterate through the slides data and create slides
    for slide_text in slides:
        slide_layout = presentation.slide_layouts[1]  # Choose the slide layout (Title and Content)
        slide = presentation.slides.add_slide(slide_layout)

        title_shape = slide.shapes.title
        title_shape.text = title

        content_slide = slide.placeholders[1]
        points = slide_text.split("\n")
        for point in points:
            content_slide.text_frame.add_paragraph().text = point

    # Save the presentation
    print("Before Save is called")
    presentation.save(r"C:\Users\dell\Desktop\Slide Generator\SlidesGenerator\templates\output.pptx")

    print("Slides created successfully!")

    return "Slides created successfully!"


if __name__ == "__main__":
    app.run(host="192.168.56.1", port=9000, debug=True)
